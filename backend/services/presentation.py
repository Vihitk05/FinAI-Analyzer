from __future__ import annotations

import hashlib
import io
import subprocess
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree

import pymupdf as fitz
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from services.export_coverage import check_coverage, coverage_failure_message
from services.presentation_model import PresentationModel, build_presentation_model


class ExportValidationError(RuntimeError):
    pass



INK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_BG = RGBColor(0x2E, 0x1F, 0x66)
TITLE_BG_ACCENT = RGBColor(0x6D, 0x28, 0xD9)
RISK_RED = RGBColor(0xEF, 0x44, 0x44)

ACCENTS = [
    RGBColor(0x63, 0x66, 0xF1), RGBColor(0xEC, 0x48, 0x99), RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0x10, 0xB9, 0x81), RGBColor(0x3B, 0x82, 0xF6), RGBColor(0x8B, 0x5C, 0xF6),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
FONT = "Calibri"


def _accent(index: int) -> RGBColor:
    return ACCENTS[index % len(ACCENTS)]


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAGE_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def _textbox(slide, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = wrap
    frame.vertical_anchor = anchor
    for i, line in enumerate(str(text).split("\n")):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT
        run.font.color.rgb = color
    return box


def _bullets(slide, x, y, w, h, items, size=14, color=INK, bullet_color=None, spacing=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    for i, item in enumerate(items):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        marker = paragraph.add_run()
        marker.text = "●  "
        marker.font.size = Pt(size)
        marker.font.name = FONT
        marker.font.color.rgb = bullet_color or color
        text_run = paragraph.add_run()
        text_run.text = str(item)
        text_run.font.size = Pt(size)
        text_run.font.name = FONT
        text_run.font.color.rgb = color
    return box


def _header_bar(slide, title: str, accent: RGBColor, subtitle: str = ""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    _textbox(slide, MARGIN, Inches(0.16), SLIDE_W - 2 * MARGIN, Inches(0.55), title.upper(), size=24, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _textbox(slide, MARGIN, Inches(0.62), SLIDE_W - 2 * MARGIN, Inches(0.35), subtitle, size=12, color=WHITE)


def _footer(slide, source_refs: list[dict]):
    if not source_refs:
        return
    seen, lines = set(), []
    for ref in source_refs[:4]:
        key = (ref.get("document"), ref.get("page"))
        if key in seen:
            continue
        seen.add(key)
        page = f" p.{ref['page']}" if ref.get("page") is not None else ""
        lines.append(f"{ref.get('document', 'Uploaded report')}{page}")
    _textbox(slide, MARGIN, SLIDE_H - Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Source: " + " | ".join(lines), size=9, color=MUTED)


def _page_number(slide, n: int):
    _textbox(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.5), Inches(0.6), Inches(0.4), str(n), size=9, color=MUTED, align=PP_ALIGN.RIGHT)





def _title_slide(prs, section, currency):
    slide = _blank_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = TITLE_BG
    band.line.fill.background()
    band.shadow.inherit = False
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(4.2), -Inches(1.6), Inches(5.5), Inches(5.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TITLE_BG_ACCENT
    accent.line.fill.background()
    accent.shadow.inherit = False
    _textbox(slide, Inches(0.7), Inches(2.9), Inches(10), Inches(1.3), section.title, size=40, color=WHITE, bold=True)
    _textbox(slide, Inches(0.7), Inches(3.9), Inches(10), Inches(0.6), "FINANCIAL INTELLIGENCE REPORT", size=18, color=ACCENTS[2], bold=True)
    if section.body.get("period"):
        _textbox(slide, Inches(0.7), Inches(4.75), Inches(10), Inches(0.5), section.body["period"], size=14, color=RGBColor(0xC7, 0xC1, 0xE8))
    _textbox(slide, Inches(0.7), SLIDE_H - Inches(0.9), Inches(11), Inches(0.5),
             "Every figure traces to a page-level source reference.", size=11, color=RGBColor(0xC7, 0xC1, 0xE8))
    return slide


def _summary_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[1])
    paragraphs = [p for p in section.body.get("paragraphs", []) if isinstance(p, str)]
    _textbox(slide, MARGIN, Inches(1.4), SLIDE_W - 2 * MARGIN, SLIDE_H - Inches(2.1),
             "\n\n".join(paragraphs), size=15, color=INK)
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _documents_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[4])
    reports = section.body.get("reports", [])
    _textbox(slide, MARGIN, Inches(1.35), SLIDE_W - 2 * MARGIN, Inches(0.5),
             f"Documents analyzed: {len(reports)}", size=16, color=INK, bold=True)
    lines = []
    for r in reports:
        period = r.get("period")
        prefix = f"{period} - " if period and period != "Undated report" else ""
        doc_type = f"  ({r.get('documentType')})" if r.get("documentType") else ""
        lines.append(f"{prefix}{r.get('filename', 'Uploaded report')}{doc_type}")
    _bullets(slide, MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(3.6), lines, size=14, bullet_color=ACCENTS[4])
    if section.body.get("period"):
        _textbox(slide, MARGIN, SLIDE_H - Inches(1.1), SLIDE_W - 2 * MARGIN, Inches(0.5),
                 f"Reporting period: {section.body['period']}", size=12, color=MUTED)
    _page_number(slide, idx)
    return slide


def _kpi_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[0])
    metrics = section.body.get("metrics", [])
    currency = section.body.get("currency", "")
    cols = min(3, len(metrics)) or 1
    gap = Inches(0.3)
    card_w = (SLIDE_W - 2 * MARGIN - gap * (cols - 1)) / cols
    card_h = Inches(1.9)
    top = Inches(1.5)
    for i, metric in enumerate(metrics):
        row, col = divmod(i, cols)
        x = MARGIN + col * (card_w + gap)
        y = top + row * (card_h + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.adjustments[0] = 0.08
        card.fill.solid()
        card.fill.fore_color.rgb = _accent(i)
        card.line.fill.background()
        card.shadow.inherit = False
        _textbox(slide, x + Inches(0.2), y + Inches(0.18), card_w - Inches(0.4), Inches(0.4),
                 str(metric.get("label", metric.get("field", ""))).upper(), size=13, color=WHITE, bold=True)
        value = metric.get("value")
        text = _money(value, currency) if isinstance(value, (int, float)) else str(value)
        _textbox(slide, x + Inches(0.2), y + Inches(0.65), card_w - Inches(0.4), Inches(1.0), text,
                 size=24, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if metric.get("status") in {"conflicting", "needs_review"}:
            _textbox(slide, x + Inches(0.2), y + card_h - Inches(0.42), card_w - Inches(0.4), Inches(0.35),
                     metric["status"].replace("_", " ").title(), size=10, color=WHITE)
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _line_chart_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[3])
    series = section.body.get("series", [])
    chart_data = CategoryChartData()
    chart_data.categories = [str(p.get("period", "")) for p in series]
    chart_data.add_series(section.body.get("unit") or "Value", tuple(p.get("value", 0) for p in series))
    frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, MARGIN, Inches(1.5), SLIDE_W - 2 * MARGIN, Inches(4.8), chart_data)
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "#,##0"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(10)
    chart.series[0].format.line.color.rgb = ACCENTS[3]
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _pie_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[4])
    breakdowns = section.body.get("breakdowns", [])
    currency = section.body.get("currency", "")
    n = max(len(breakdowns), 1)
    chart_w = (SLIDE_W - 2 * MARGIN) / n
    for i, bd in enumerate(breakdowns):
        rows = [r for r in bd.get("rows", []) if isinstance(r, dict)][:8]
        x = MARGIN + i * chart_w
        _textbox(slide, x, Inches(1.25), chart_w - Inches(0.2), Inches(0.4), str(bd.get("subtitle", "")).upper(), size=13, color=MUTED, bold=True)
        chart_data = CategoryChartData()
        chart_data.categories = [str(r.get("category", "")) for r in rows]
        chart_data.add_series(bd.get("subtitle", ""), tuple((r.get("amount") or 0) for r in rows))
        frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, Inches(1.7), chart_w - Inches(0.3), Inches(4.6), chart_data)
        chart = frame.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = "#,##0"
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.font.size = Pt(9)
        for j, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _accent(j)
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _ratio_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[0])
    rows = section.body.get("rows", [])
    per_slide = 12
    _rows_block(slide, rows[:per_slide])
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _rows_block(slide, rows):
    top = Inches(1.45)
    row_h = Inches(0.46)
    for i, row in enumerate(rows):
        y = top + i * row_h
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, Inches(0.08), row_h - Inches(0.1))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _accent(i)
        stripe.line.fill.background()
        stripe.shadow.inherit = False
        _textbox(slide, MARGIN + Inches(0.25), y, Inches(8.5), row_h, str(row.get("ratio", "")), size=13, color=INK)
        value = row.get("value")
        text = f"{value}{'%' if row.get('isPercentage') else ''}" if value is not None else ""
        _textbox(slide, SLIDE_W - MARGIN - Inches(2.6), y, Inches(2.4), row_h, text, size=13, color=ACCENTS[0], bold=True, align=PP_ALIGN.RIGHT)


def _table_slide(prs, section, idx, *, columns=None, rows=None, title=None):
    columns = columns or section.body.get("columns", [])
    rows = rows if rows is not None else section.body.get("rows", [])
    slide = _blank_slide(prs)
    _header_bar(slide, title or section.title, ACCENTS[5])
    if not columns:
        _textbox(slide, MARGIN, Inches(1.5), SLIDE_W - 2 * MARGIN, Inches(1), "No data", size=14, color=MUTED)
        _page_number(slide, idx)
        return slide
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, len(columns), MARGIN, Inches(1.4),
                                         SLIDE_W - 2 * MARGIN, Inches(min(5.4, 0.42 * n_rows + 0.4))).table
    for c, name in enumerate(columns):
        cell = table_shape.cell(0, c)
        cell.text = str(name)
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENTS[5]
    for r, row in enumerate(rows, start=1):
        for c in range(len(columns)):
            cell = table_shape.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            run = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else cell.text_frame.paragraphs[0].add_run()
            run.font.size = Pt(10)
            run.font.color.rgb = INK
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _narrative_slide(prs, section, idx):
    slide = _blank_slide(prs)
    accent = RISK_RED if section.body.get("accent") == "risk" else ACCENTS[2]
    _header_bar(slide, section.title, accent)
    bullets = [str(b) for b in section.body.get("bullets", []) if b][:12]
    _bullets(slide, MARGIN, Inches(1.4), SLIDE_W - 2 * MARGIN, SLIDE_H - Inches(2.1), bullets, size=15, color=INK, bullet_color=accent)
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _notes_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[0])
    y = Inches(1.35)
    for note in section.body.get("notes", []):
        _textbox(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.4), note.get("category", "").upper(), size=13, color=ACCENTS[0], bold=True)
        y += Inches(0.42)
        items = [str(i) for i in note.get("items", [])][:5]
        _bullets(slide, MARGIN + Inches(0.15), y, SLIDE_W - 2 * MARGIN - Inches(0.15), Inches(0.34 * len(items) + 0.1), items, size=12)
        y += Inches(0.34 * len(items) + 0.25)
    _footer(slide, section.source_refs)
    _page_number(slide, idx)
    return slide


def _data_status_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, RISK_RED)
    y = Inches(1.4)
    conflicts = section.body.get("conflicts", [])
    if conflicts:
        _textbox(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.4), "SOURCE CONFLICTS", size=13, color=RISK_RED, bold=True)
        y += Inches(0.45)
        lines = [f"{c.get('field')}: sources disagree by ~{c.get('spreadPercent')}% for the same period (shown unresolved)" for c in conflicts]
        _bullets(slide, MARGIN + Inches(0.15), y, SLIDE_W - 2 * MARGIN, Inches(0.34 * len(lines) + 0.1), lines, size=12, bullet_color=RISK_RED)
        y += Inches(0.34 * len(lines) + 0.35)
    missing = section.body.get("missing", [])
    if missing:
        _textbox(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.4), "NOT IDENTIFIED IN THE VERIFIED SOURCES", size=13, color=MUTED, bold=True)
        y += Inches(0.45)
        _bullets(slide, MARGIN + Inches(0.15), y, SLIDE_W - 2 * MARGIN, Inches(0.34 * len(missing) + 0.1),
                 [m.get("message", m.get("label", "")) for m in missing], size=12, color=MUTED)
    _page_number(slide, idx)
    return slide


def _sources_slide(prs, section, idx):
    slide = _blank_slide(prs)
    _header_bar(slide, section.title, ACCENTS[4])
    lines = []
    for r in section.body.get("reports", []):
        period = f" - {r.get('period')}" if r.get("period") and r.get("period") != "Undated report" else ""
        doc_type = f" ({r.get('documentType')})" if r.get("documentType") else ""
        lines.append(f"{r.get('filename', 'Uploaded report')}{period}{doc_type}")
    _bullets(slide, MARGIN, Inches(1.45), SLIDE_W - 2 * MARGIN, Inches(3.6), lines, size=14, bullet_color=ACCENTS[4])
    external = section.body.get("external", {})
    if external.get("message"):
        _textbox(slide, MARGIN, SLIDE_H - Inches(1.0), SLIDE_W - 2 * MARGIN, Inches(0.5), external["message"], size=10, color=MUTED)
    _page_number(slide, idx)
    return slide


_KIND_BUILDERS = {
    "summary": _summary_slide,
    "documents": _documents_slide,
    "kpi": _kpi_slide,
    "line_chart": _line_chart_slide,
    "pie": _pie_slide,
    "ratios": _ratio_slide,
    "table": _table_slide,
    "narrative": _narrative_slide,
    "notes": _notes_slide,
    "data_status": _data_status_slide,
    "sources": _sources_slide,
}


def _money(value, currency: str) -> str:
    if not isinstance(value, (int, float)):
        return str(value)
    a = abs(value)
    for scale, suffix in ((1e12, "T"), (1e9, "B"), (1e6, "M"), (1e3, "K")):
        if a >= scale:
            return f"{currency} {value / scale:,.2f}{suffix}".strip()
    return f"{currency} {value:,.0f}".strip()





def build_pptx(data: dict, model: PresentationModel | None = None) -> bytes:
    model = model or build_presentation_model(data)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, section in enumerate(model.sections, start=1):
        if section.kind == "title":
            _title_slide(prs, section, model.currency)
            continue
        builder = _KIND_BUILDERS.get(section.kind)
        if builder is None:
            continue

        if section.kind == "table" and len(section.body.get("rows", [])) > 14:
            rows = section.body["rows"]
            for part, start in enumerate(range(0, len(rows), 14)):
                _table_slide(prs, section, idx, columns=section.body["columns"], rows=rows[start:start + 14],
                             title=f"{section.title} ({part + 1})" if len(rows) > 14 else section.title)
        else:
            builder(prs, section, idx)

    buffer = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    path = Path(buffer.name)
    buffer.close()
    try:
        prs.save(str(path))
        return path.read_bytes()
    finally:
        path.unlink(missing_ok=True)




_REQUIRED_OPC_PARTS = ("[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml", "ppt/_rels/presentation.xml.rels")


def _validate_opc_package(pptx_bytes: bytes) -> int:

    try:
        archive = zipfile.ZipFile(io.BytesIO(pptx_bytes))
    except zipfile.BadZipFile as exc:
        raise ExportValidationError(f"Generated PPTX is not a valid zip: {exc}") from exc
    if archive.testzip() is not None:
        raise ExportValidationError("Generated PPTX zip has a corrupt member")
    names = set(archive.namelist())
    missing = [p for p in _REQUIRED_OPC_PARTS if p not in names]
    if missing:
        raise ExportValidationError(f"Generated PPTX is missing required package parts: {', '.join(missing)}")

    slide_parts = sorted(n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    if not slide_parts:
        raise ExportValidationError("Generated PPTX contains no slide parts")
    for part in slide_parts:
        try:
            ElementTree.fromstring(archive.read(part))
        except ElementTree.ParseError as exc:
            raise ExportValidationError(f"Slide part {part} is not well-formed XML: {exc}") from exc
        rels = f"ppt/slides/_rels/{part.split('/')[-1]}.rels"
        if rels not in names:
            raise ExportValidationError(f"Slide part {part} has no relationships part")

    try:
        pres_xml = ElementTree.fromstring(archive.read("ppt/presentation.xml"))
    except ElementTree.ParseError as exc:
        raise ExportValidationError(f"ppt/presentation.xml is not well-formed: {exc}") from exc
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    listed = pres_xml.findall(".//p:sldIdLst/p:sldId", ns)
    if len(listed) != len(slide_parts):
        raise ExportValidationError(
            f"presentation.xml lists {len(listed)} slides but the package has {len(slide_parts)} slide parts"
        )
    return len(slide_parts)


def validate_pptx(pptx_bytes: bytes) -> int:

    opc_slides = _validate_opc_package(pptx_bytes)

    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        raise ExportValidationError(f"Generated PPTX could not be opened: {exc}") from exc

    slide_count = len(prs.slides)
    if slide_count == 0:
        raise ExportValidationError("Generated PPTX has no slides")
    if slide_count != opc_slides:
        raise ExportValidationError(f"python-pptx sees {slide_count} slides, package has {opc_slides}")

    for index, slide in enumerate(prs.slides, start=1):
        has_content = any(
            (getattr(shape, "has_chart", False) and shape.has_chart)
            or (getattr(shape, "has_table", False) and shape.has_table)
            or (getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())
            for shape in slide.shapes
        )
        if not has_content:
            raise ExportValidationError(f"Slide {index} has no readable content")

    return slide_count


def enforce_coverage(data: dict, model: PresentationModel) -> dict:

    from services.logging_config import get_logger

    coverage = check_coverage(data, model)
    structural_gap = bool(coverage["missingCharts"] or coverage["missingTables"] or len(coverage["missingKeys"]) > 2)
    if structural_gap:
        raise ExportValidationError(coverage_failure_message(coverage))
    if coverage["missingKeys"]:
        get_logger(__name__).warning(
            "export_coverage_incomplete", extra={"missing": coverage["missingKeys"]}
        )
    return coverage





def pdf_export(pptx_bytes: bytes, data: dict | None = None, model: PresentationModel | None = None) -> tuple[bytes, str]:

    try:
        return _pdf_via_soffice(pptx_bytes), "soffice"
    except (OSError, subprocess.SubprocessError, RuntimeError) as primary_exc:
        if data is None and model is None:
            raise RuntimeError(f"PDF export is unavailable on this server: {primary_exc}") from primary_exc
        try:
            return _pdf_via_pymupdf(model or build_presentation_model(data)), "fallback"
        except Exception as fallback_exc:
            raise RuntimeError(
                f"PDF export failed on both paths (primary: {primary_exc}; fallback: {fallback_exc})"
            ) from fallback_exc


def _pdf_via_soffice(pptx_bytes: bytes) -> bytes:
    with tempfile.TemporaryDirectory() as directory:
        pptx = Path(directory) / "company-intelligence.pptx"
        pptx.write_bytes(pptx_bytes)
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", directory, str(pptx)],
            check=True, capture_output=True, timeout=60,
        )
        pdf = pptx.with_suffix(".pdf")
        if not pdf.exists():
            raise RuntimeError("LibreOffice did not produce a PDF file")
        return pdf.read_bytes()


_PDF_W, _PDF_H = 792, 612


def _pdf_via_pymupdf(model: PresentationModel) -> bytes:
    doc = fitz.open()

    def new_page(title: str):
        page = doc.new_page(width=_PDF_W, height=_PDF_H)
        page.draw_rect(fitz.Rect(0, 0, _PDF_W, 60), color=None, fill=(0.18, 0.12, 0.4))
        page.insert_text((40, 38), title.upper(), fontsize=20, color=(1, 1, 1), fontname="hebo")
        return page

    def body_lines(page, lines, start_y=90, size=11, color=(0.12, 0.16, 0.23)):
        y = start_y
        for line in lines:
            if y > _PDF_H - 50:
                page = new_page("(continued)")
                y = 90
            page.insert_text((44, y), str(line)[:150], fontsize=size, color=color)
            y += size + 7
        return page

    for section in model.sections:
        if section.kind == "title":
            page = doc.new_page(width=_PDF_W, height=_PDF_H)
            page.draw_rect(fitz.Rect(0, 0, _PDF_W, _PDF_H), fill=(0.18, 0.12, 0.4))
            page.insert_text((60, 300), section.title, fontsize=34, color=(1, 1, 1), fontname="hebo")
            page.insert_text((60, 340), "FINANCIAL INTELLIGENCE REPORT", fontsize=14, color=(0.96, 0.62, 0.04))
            if section.body.get("period"):
                page.insert_text((60, 372), str(section.body["period"]), fontsize=12, color=(0.78, 0.76, 0.91))
            continue

        page = new_page(section.title)
        b = section.body
        if section.kind == "kpi":
            currency = b.get("currency", "")
            page = body_lines(page, [f"{m.get('label', m.get('field'))}:  {_money(m.get('value'), currency)}" for m in b.get("metrics", [])], size=13)
        elif section.kind in {"summary", "narrative"}:
            text = b.get("paragraphs") or b.get("bullets") or []
            wrapped = []
            for para in text:
                s = str(para)
                while s:
                    wrapped.append(("• " if section.kind == "narrative" else "") + s[:120])
                    s = s[120:]
            page = body_lines(page, wrapped)
        elif section.kind in {"table", "ratios"}:
            if section.kind == "ratios":
                cols = ["Ratio", "Value"]
                rows = [[r.get("ratio", ""), f"{r.get('value')}{'%' if r.get('isPercentage') else ''}"] for r in b.get("rows", [])]
            else:
                cols, rows = b.get("columns", []), b.get("rows", [])
            page = body_lines(page, ["   ".join(f"{c:>14}" for c in cols)] + ["   ".join(f"{str(c):>14}" for c in row) for row in rows], size=10)
        elif section.kind == "line_chart":
            page = _pdf_bar(page, [(str(p.get("period", "")), p.get("value", 0)) for p in b.get("series", [])], b.get("currency", ""))
        elif section.kind == "pie":
            lines = []
            for bd in b.get("breakdowns", []):
                lines.append(bd.get("subtitle", "").upper())
                for r in bd.get("rows", [])[:8]:
                    lines.append(f"   {r.get('category', '')}: {_money(r.get('amount'), b.get('currency', ''))}")
            page = body_lines(page, lines)
        elif section.kind == "documents":
            page = body_lines(page, [f"Documents analyzed: {len(b.get('reports', []))}"] +
                              [f"  • {r.get('period', '')} - {r.get('filename', '')}" for r in b.get("reports", [])])
        elif section.kind == "notes":
            lines = []
            for note in b.get("notes", []):
                lines.append(note.get("category", "").upper())
                lines.extend(f"  • {i}" for i in note.get("items", []))
            page = body_lines(page, lines)
        elif section.kind == "data_status":
            lines = [f"CONFLICT - {c.get('field')} (~{c.get('spreadPercent')}% spread)" for c in b.get("conflicts", [])]
            lines += [f"MISSING - {m.get('message', m.get('label', ''))}" for m in b.get("missing", [])]
            page = body_lines(page, lines)
        elif section.kind == "sources":
            page = body_lines(page, [f"• {r.get('filename', '')} - {r.get('period', '')}" for r in b.get("reports", [])])

        if section.source_refs:
            refs = " | ".join(f"{r.get('document', '')} p.{r.get('page')}" for r in section.source_refs[:3] if r.get("page") is not None)
            page.insert_text((44, _PDF_H - 30), f"Source: {refs}", fontsize=8, color=(0.39, 0.45, 0.55))

    out = doc.tobytes()
    doc.close()
    return out


def _pdf_bar(page, points, currency):
    if not points:
        page.insert_text((44, 100), "No data", fontsize=11, color=(0.39, 0.45, 0.55))
        return page
    x0, y0, w, h = 60, _PDF_H - 80, _PDF_W - 140, 380
    peak = max((abs(v) for _, v in points), default=1) or 1
    bar_w = w / (len(points) * 1.5)
    for i, (label_text, value) in enumerate(points):
        bx = x0 + i * (w / len(points))
        bh = (abs(value) / peak) * h
        page.draw_rect(fitz.Rect(bx, y0 - bh, bx + bar_w, y0), fill=(0.06, 0.72, 0.51))
        page.insert_text((bx, y0 + 14), str(label_text)[:12], fontsize=8, color=(0.28, 0.33, 0.41))
        page.insert_text((bx, y0 - bh - 6), _money(value, currency), fontsize=8, color=(0.12, 0.16, 0.23))
    return page


def validate_pdf(pdf_bytes: bytes, expected_page_count: int, *, exact: bool = True) -> int:

    try:
        document = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as exc:
        raise ExportValidationError(f"Generated PDF could not be parsed: {exc}") from exc
    try:
        page_count = document.page_count
        if page_count == 0:
            raise ExportValidationError("Generated PDF has no pages")
        if exact and page_count != expected_page_count:
            raise ExportValidationError(
                f"PDF page count ({page_count}) does not match the presentation's slide count ({expected_page_count})"
            )
        if not exact and page_count < expected_page_count:
            raise ExportValidationError(
                f"PDF has fewer pages ({page_count}) than the presentation has slides ({expected_page_count})"
            )
        full_text = ""
        empty_pages = []
        for i in range(page_count):
            text = document[i].get_text().strip()
            if not text:
                empty_pages.append(i + 1)
            full_text += text + "\n"
        if empty_pages:
            raise ExportValidationError(f"PDF page(s) {empty_pages} converted with no readable content")
        for marker in ('{"field"', '"citations":', '"reportId"', "custom_id"):
            if marker in full_text:
                raise ExportValidationError(f"PDF contains raw data ({marker!r}) that must not reach a stakeholder file")
        return page_count
    finally:
        document.close()


def file_checksum(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()
